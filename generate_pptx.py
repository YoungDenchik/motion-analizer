"""
Generate AI Fitness Coach PPTX presentation.
Run: python generate_pptx.py
Output: AI_Fitness_Coach.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Palette ──────────────────────────────────────────────────────────────────
BG      = RGBColor(0x0D, 0x14, 0x2B)
BAR     = RGBColor(0x08, 0x12, 0x26)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
BLUE    = RGBColor(0x42, 0xA5, 0xF5)
LBLUE   = RGBColor(0x90, 0xCA, 0xF9)
GRAY    = RGBColor(0xCF, 0xD8, 0xDC)
DIM     = RGBColor(0x78, 0x90, 0xA0)
CODEBG  = RGBColor(0x07, 0x0C, 0x1A)
CODEFG  = RGBColor(0xA8, 0xD0, 0xF0)
GREEN   = RGBColor(0x66, 0xBB, 0x6A)
YELLOW  = RGBColor(0xFF, 0xD5, 0x4F)
RED     = RGBColor(0xEF, 0x53, 0x50)
TH      = RGBColor(0x1A, 0x3C, 0x72)
TR1     = RGBColor(0x0F, 0x20, 0x44)
TR2     = RGBColor(0x0A, 0x18, 0x38)
BOXBG   = RGBColor(0x10, 0x1E, 0x3C)

# Slide dimensions 16:9
W = Inches(13.333)
H = Inches(7.5)

prs = Presentation()
prs.slide_width = W
prs.slide_height = H
BLANK = prs.slide_layouts[6]


# ── Low-level helpers ─────────────────────────────────────────────────────────

def new_slide():
    s = prs.slides.add_slide(BLANK)
    f = s.background.fill
    f.solid()
    f.fore_color.rgb = BG
    return s


def box(s, l, t, w, h, fill=None, line=None, lw=0.75):
    sh = s.shapes.add_shape(1, l, t, w, h)
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.color.rgb = line
        sh.line.width = Pt(lw)
    else:
        sh.line.fill.background()
    return sh


def label(s, text, l, t, w, h, sz=18, bold=False, italic=False,
          fg=GRAY, align=PP_ALIGN.LEFT, mono=False):
    b = s.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = fg
    r.font.name = "Consolas" if mono else "Calibri"
    return b, tf


def multiline(s, lines, l, t, w, h, def_sz=17, def_fg=GRAY, mono=False):
    """
    Multi-paragraph text box.
    Each line: str  or  dict(text, sz, fg, bold, italic, mono, align, sp_b, sp_a)
    """
    b = s.shapes.add_textbox(l, t, w, h)
    tf = b.text_frame
    tf.word_wrap = True
    first = True
    for ln in lines:
        if isinstance(ln, str):
            d = {'text': ln}
        else:
            d = ln
        text = d.get('text', '')
        sz   = d.get('sz', def_sz)
        fg   = d.get('fg', def_fg)
        bold = d.get('bold', False)
        ital = d.get('italic', False)
        mn   = d.get('mono', mono)
        aln  = d.get('align', PP_ALIGN.LEFT)
        sp_b = d.get('sp_b', 0)
        sp_a = d.get('sp_a', 0)

        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = aln
        if sp_b:
            p.space_before = Pt(sp_b)
        if sp_a:
            p.space_after = Pt(sp_a)

        r = p.add_run()
        r.text = text
        r.font.size = Pt(sz)
        r.font.bold = bold
        r.font.italic = ital
        r.font.color.rgb = fg
        r.font.name = "Consolas" if mn else "Calibri"
    return b


def top_bar(s, title, sub=None):
    box(s, 0, 0, W, Inches(1.15), fill=BAR)
    box(s, 0, Inches(1.15), W, Inches(0.07), fill=BLUE)
    label(s, title, Inches(0.45), Inches(0.1), W - Inches(0.9), Inches(0.82),
          sz=30, bold=True, fg=WHITE)
    if sub:
        label(s, sub, Inches(0.45), Inches(0.82), W - Inches(0.9), Inches(0.36),
              sz=15, fg=LBLUE)


def code_block(s, text, l, t, w, h, sz=12.5):
    box(s, l, t, w, h, fill=CODEBG, line=BLUE)
    pad = Inches(0.18)
    label(s, text, l + pad, t + pad, w - 2 * pad, h - 2 * pad,
          sz=sz, fg=CODEFG, mono=True)


def mk_table(s, hdrs, rows, l, t, w, h, col_widths=None, hsz=14, rsz=13,
             center_all=False):
    nc = len(hdrs)
    nr = len(rows) + 1
    tbl = s.shapes.add_table(nr, nc, l, t, w, h).table

    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = int(cw)
    else:
        cw = w // nc
        for ci in range(nc):
            tbl.columns[ci].width = cw

    for ci, hdr in enumerate(hdrs):
        c = tbl.cell(0, ci)
        c.fill.solid()
        c.fill.fore_color.rgb = TH
        p = c.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = hdr
        r.font.size = Pt(hsz)
        r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "Calibri"

    for ri, row in enumerate(rows):
        bg = TR1 if ri % 2 == 0 else TR2
        for ci, val in enumerate(row):
            c = tbl.cell(ri + 1, ci)
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            p = c.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if center_all else (
                PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER)
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(rsz)
            r.font.color.rgb = GRAY
            r.font.name = "Calibri"
    return tbl


# Content Y start (below title bar + accent)
CY = Inches(1.3)
CW = W - Inches(0.9)  # usable content width
CH = H - Inches(1.45)  # usable content height
CL = Inches(0.45)     # content left margin

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
box(s, 0, H - Inches(0.55), W, Inches(0.55), fill=BAR)
box(s, 0, H - Inches(0.55), W, Inches(0.07), fill=BLUE)

label(s, "AI Fitness Coach",
      Inches(0.8), Inches(1.2), Inches(11.8), Inches(2.0),
      sz=64, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)

label(s, "Розумний персональний тренер у смартфоні",
      Inches(0.8), Inches(3.0), Inches(11.8), Inches(0.75),
      sz=28, fg=LBLUE, align=PP_ALIGN.CENTER)

label(s, "Аналіз техніки виконання вправ через комп'ютерний зір\n"
          "та динамічне порівняння з еталонним рухом",
      Inches(1.5), Inches(3.8), Inches(10.4), Inches(1.0),
      sz=20, fg=GRAY, align=PP_ALIGN.CENTER)

box(s, Inches(5.2), Inches(5.2), Inches(3.0), Inches(0.6), fill=TH, line=BLUE)
label(s, "Прототип  v0.1",
      Inches(5.2), Inches(5.2), Inches(3.0), Inches(0.6),
      sz=18, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Проблема")

multiline(s, [
    {'text': '  80% людей займаються без тренера — і роблять це неправильно',
     'sz': 21, 'fg': WHITE, 'bold': True, 'sp_a': 4},
    {'text': '  Неправильна техніка → мікротравми → хронічні ушкодження суглобів',
     'sz': 19, 'fg': GRAY, 'sp_a': 4},
    {'text': '  Персональний тренер: дорого, недоступно 24/7',
     'sz': 19, 'fg': GRAY, 'sp_a': 4},
    {'text': '  YouTube-відео не дають персонального фідбеку',
     'sz': 19, 'fg': GRAY, 'sp_a': 4},
    {'text': '  Спортзали не мають інструменту масштабованого контролю техніки',
     'sz': 19, 'fg': GRAY},
], CL, CY, CW, Inches(3.6))

box(s, CL, Inches(5.25), CW, Inches(1.75), fill=TR1, line=BLUE)
label(s, '"Я думав, що роблю все правильно —\nпоки не отримав травму коліна."',
      CL + Inches(0.3), Inches(5.4), CW - Inches(0.6), Inches(1.45),
      sz=21, italic=True, fg=LBLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Рішення: AI Fitness Coach",
        "Система, що бачить вправу і коригує виконання")

steps = [
    ("Відео користувача", BLUE),
    ("Детекція 33 точок скелету  (MediaPipe)", GRAY),
    ("9 кутових характеристик суглобів", GRAY),
    ("Порівняння з еталоном через DTW", GRAY),
    ("Оцінка 0–100 + класифікація помилок", GRAY),
    ("Персональний тренерський фідбек", GREEN),
]

sx, sy = Inches(0.6), Inches(1.45)
sw, sh = Inches(12.1), Inches(0.72)
gap = Inches(0.12)

for i, (text, color) in enumerate(steps):
    y = sy + i * (sh + gap)
    bg = TH if i == 0 else BOXBG
    box(s, sx, y, sw, sh, fill=bg, line=BLUE)
    arrow_y = y + sh
    if i < len(steps) - 1:
        box(s, Inches(6.4), arrow_y, Inches(0.55), gap, fill=BLUE)
    label(s, text, sx + Inches(0.3), y, sw - Inches(0.6), sh,
          sz=19, bold=(i == 0), fg=color)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Business model
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Цільова аудиторія та бізнес-модель")

cols = [
    ("Спортзали  (B2B)",
     ["Інтеграція в інфраструктуру залу",
      "Аналіз техніки відвідувачів",
      "Звіти для менеджменту",
      "Монетизація через підписку залу"]),
    ("Тренери  (B2B2C)",
     ["Запис власного еталону вправи",
      "Аналіз клієнтів онлайн та офлайн",
      "Масштабування без ручної роботи",
      "Платформа для онлайн-коучингу"]),
    ("Кінцевий користувач  (B2C)",
     ["Мобільний додаток",
      "Самостійні тренування",
      "Прогрес і аналітика",
      "Персональний план тренувань"]),
]

cw = Inches(4.0)
gap = Inches(0.27)
for i, (title, bullets) in enumerate(cols):
    cx = CL + i * (cw + gap)
    box(s, cx, CY, cw, Inches(5.7), fill=BOXBG, line=BLUE)
    label(s, title, cx + Inches(0.2), CY + Inches(0.15), cw - Inches(0.4), Inches(0.6),
          sz=17, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    box(s, cx, CY + Inches(0.75), cw, Inches(0.05), fill=BLUE)
    lines = [{'text': f"  •  {b}", 'sz': 16, 'fg': GRAY, 'sp_a': 6} for b in bullets]
    multiline(s, lines, cx + Inches(0.2), CY + Inches(0.9),
              cw - Inches(0.4), Inches(4.6))

label(s, "Тренер записує власний еталон вправи → система навчається під його методику",
      CL, Inches(7.05), CW, Inches(0.4), sz=15, fg=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — What's implemented
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Що реалізовано у прототипі")

rows_impl = [
    ("Детекція пози — MediaPipe 33 landmarks/кадр",    "✅ Готово"),
    ("9 кутових ознак суглобів + нормалізація",        "✅ Готово"),
    ("Сегментація повторів (PCA + multi-scale peaks)", "✅ Готово"),
    ("FastDTW порівняння з ваговими коефіцієнтами",   "✅ Готово"),
    ("Толерантна зона референсу (міжповторна σ)",     "✅ Готово"),
    ("Оцінка 0–100 + літерний грейд (A–F)",           "✅ Готово"),
    ("Класифікація: Critical (>12°) / Technical (>8°)","✅ Готово"),
    ("Генерація тренерських реплік (rule-based)",     "✅ Готово"),
    ("Анотація відео (скелет, шкала, фідбек-банер)",  "✅ Готово"),
    ("Motion Signature 99-dim + автовибір референсу", "✅ Готово"),
    ("CLI: record / analyze / list / demo",           "✅ Готово"),
]

mk_table(s, ["Компонент", "Статус"], rows_impl,
         CL, CY, CW, Inches(5.85),
         col_widths=[CW * 0.80, CW * 0.20],
         hsz=15, rsz=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — Architecture
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Загальна архітектура системи")

arch = (
    "  ВІДЕО ДЖЕРЕЛО:  Файл  |  Камера  |  Стрім\n"
    "           │\n"
    "  ┌────────▼────────┐\n"
    "  │  PoseDetector   │  MediaPipe — 33 landmarks (x, y, z, vis)\n"
    "  └────────┬────────┘\n"
    "           │\n"
    "  ┌────────▼────────┐\n"
    "  │FeatureExtractor │  9 кутових ознак суглобів + нормалізація\n"
    "  └────────┬────────┘\n"
    "           │\n"
    "  ┌────────▼────────┐\n"
    "  │   TimeSeries    │  (T × 9) float32  +  Savitzky-Golay згладжування\n"
    "  └───────┬─────────┘\n"
    "          │\n"
    "   ┌──────┴──────┐\n"
    "   │             │\n"
    "RepSegmenter  ReferenceStore → find_best_match (cosine sim)\n"
    "   │             │\n"
    "   └──────┬──────┘\n"
    "          │\n"
    "  ┌───────▼──────────┐\n"
    "  │  RepComparator   │  per-rep DTW  +  tolerance band\n"
    "  └──┬───────┬────┬──┘\n"
    "     │       │    │\n"
    " Scoring  Error  Feedback → VideoRenderer"
)
code_block(s, arch, CL, CY, CW, Inches(5.85), sz=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Pipeline step-by-step
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Конвеєр обробки: крок за кроком")

pipe = (
    "Відео\n"
    "  └─ PoseDetector ──────────  33 landmarks / кадр (MediaPipe Tasks API)\n"
    "       └─ FeatureExtractor ──  9 кутів суглобів / кадр\n"
    "            └─ TimeSeries ───  (T × 9) згладжений масив\n"
    "                 └─ RepSegmenter ─── PCA-сигнал + multi-scale peaks\n"
    "                      └─ RepComparator\n"
    "                           ├─ DTWComparator ──── вирівнювання по часу\n"
    "                           ├─ ScoringEngine ──── оцінка 0–100\n"
    "                           ├─ ErrorClassifier ── класифікація помилок\n"
    "                           └─ FeedbackGenerator  тренерські репліки\n"
    "\n"
    "Відеоанотація (окремий другий прохід по вихідному відео):\n"
    "  AnalysisResult → VideoRenderer → FrameAnnotator → .mp4"
)
code_block(s, pipe, CL, CY, CW, Inches(4.4), sz=13)

multiline(s, [
    {'text': 'Чому другий прохід?', 'sz': 17, 'fg': WHITE, 'bold': True},
    {'text': 'Аналіз зберігає кути суглобів (float32), а не піксельні координати.', 'sz': 16, 'fg': GRAY},
    {'text': 'Для малювання скелету детектор запускається повторно у static_image_mode.', 'sz': 16, 'fg': GRAY},
], CL, Inches(5.9), CW, Inches(1.45))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Key technical decisions
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Ключові технічні рішення")

decisions = [
    ("Кути суглобів, а не координати",
     "Координати залежать від відстані до камери, зросту, ракурсу.\n"
     "Кут суглоба — чистий біомеханічний показник: 90° коліна = 90° завжди.\n"
     "Нормалізація: hip_center + torso_length → інваріант до позиції.",
     BLUE),
    ("FastDTW замість Euclidean",
     "Euclidean потребує однакову довжину і часове вирівнювання.\n"
     "DTW знаходить оптимальне вирівнювання по часовій осі.\n"
     "FastDTW: O(N) — реальний час для до 500 кадрів.",
     YELLOW),
    ("PCA для сегментації повторів",
     "PCA (перша головна компонента) об'єднує рухи всіх 9 суглобів.\n"
     "Стійкіший сигнал, ніж будь-який один суглоб окремо.\n"
     "Multi-scale smoothing прибирає внутрішньоповторний шум.",
     GREEN),
]

bw = Inches(4.0)
gap = Inches(0.27)
for i, (title, body, color) in enumerate(decisions):
    bx = CL + i * (bw + gap)
    box(s, bx, CY, bw, Inches(5.7), fill=BOXBG, line=color, lw=1.2)
    box(s, bx, CY, bw, Inches(0.65), fill=TH)
    label(s, title, bx + Inches(0.18), CY + Inches(0.08), bw - Inches(0.36), Inches(0.55),
          sz=16, bold=True, fg=color)
    multiline(s, [{'text': body, 'sz': 15, 'fg': GRAY}],
              bx + Inches(0.18), CY + Inches(0.8), bw - Inches(0.36), Inches(4.7))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Recording reference
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Додавання референсного прикладу")

label(s, "CLI-команди", CL, CY, CW, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s,
    "# Один або декілька повторів у одному відео\n"
    "python main.py record squat --video reference_squat.mp4\n\n"
    "# Варіанти для різних кутів зйомки / технік\n"
    "python main.py record squat_lowbar  --video lowbar_ref.mp4\n"
    "python main.py record squat_front   --video front_angle.mp4",
    CL, CY + Inches(0.5), CW, Inches(1.9), sz=13)

label(s, "Що відбувається всередині", CL, CY + Inches(2.55), CW, Inches(0.45),
      sz=17, bold=True, fg=LBLUE)
code_block(s,
    "Відео → PoseDetector → FeatureExtractor → TimeSeries\n"
    "     ↓\n"
    "RepSegmenter: знаходить границі повторів\n"
    "     ↓  якщо N повторів\n"
    "Усереднення в 1 canonical rep (RepComparator._build_canonical_rep)\n"
    "     ↓\n"
    "MotionSignature.from_timeseries()  →  99-dim vector\n"
    "     ↓\n"
    "Збереження JSON:  time_series (T×9) + signature (99 float32) + metadata",
    CL, CY + Inches(3.05), CW, Inches(2.9), sz=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Streaming
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Підтримка потокового відео")

multiline(s, [
    {'text': 'Архітектура підтримує два режими входу:', 'sz': 19, 'fg': WHITE, 'bold': True, 'sp_a': 6},
    {'text': '  •  Файл (офлайн-аналіз) — відео після тренування', 'sz': 17, 'fg': GRAY, 'sp_a': 4},
    {'text': '  •  Камера (онлайн / стрімінг) — відео в реальному часі', 'sz': 17, 'fg': GRAY},
], CL, CY, CW, Inches(1.6))

code_block(s,
    "# Відеофайл\n"
    "python main.py analyze squat --video workout.mp4 --output result.mp4\n\n"
    "# Камера (0 = перша вебкамера)\n"
    "python main.py analyze squat --video 0",
    CL, CY + Inches(1.7), CW, Inches(1.75), sz=14)

label(s, "MediaPipe Running Modes:", CL, CY + Inches(3.6), CW, Inches(0.45),
      sz=17, bold=True, fg=LBLUE)
code_block(s,
    "RunningMode.VIDEO       → синхронна обробка кадрів (файл / запис)\n"
    "RunningMode.LIVE_STREAM → асинхронна з callback   (камера / стрім)\n\n"
    "Мобільний SDK (Android / iOS): той самий конвеєр прямо на пристрої",
    CL, CY + Inches(4.1), CW, Inches(1.75), sz=14)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Storage format
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Зберігання референсів")

lw = Inches(6.1)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

label(s, "references/squat.json", CL, CY, lw, Inches(0.42),
      sz=16, bold=True, fg=LBLUE, mono=True)
code_block(s,
    '{\n'
    '  "name": "squat",\n'
    '  "description": "Side view, barbell back squat",\n'
    '  "time_series": {\n'
    '    "fps": 30.0,\n'
    '    "feature_names": ["left_knee_angle", ...],\n'
    '    "data": [[171.2, 170.8, ...], ...]\n'
    '  },\n'
    '  "signature": {\n'
    '    "vector": [131.4, 128.7, ..., 0.82],\n'
    '    "feature_names": ["left_knee_angle", ...]\n'
    '  },\n'
    '  "metadata": { "author": "coach", "date": "2025-05" }\n'
    '}',
    CL, CY + Inches(0.48), lw, Inches(3.7), sz=12)

label(s, "Чому JSON, а не HDF5?", rx, CY, rw, Inches(0.42),
      sz=16, bold=True, fg=LBLUE)
mk_table(s,
    ["Критерій", "JSON", "HDF5"],
    [("Читабельність", "✅", "❌"),
     ("Git-friendly",  "✅", "❌"),
     ("Без залежностей","✅","❌"),
     ("Швидкість",     "OK", "Краще"),
     ("<5 MB файли",   "✅", "—")],
    rx, CY + Inches(0.48), rw, Inches(2.95),
    col_widths=[rw*0.50, rw*0.25, rw*0.25],
    hsz=13, rsz=12, center_all=False)

multiline(s, [
    {'text': 'Зворотна сумісність:', 'sz': 15, 'fg': WHITE, 'bold': True},
    {'text': 'JSON без поля "signature" → перераховується автоматично при завантаженні', 'sz': 14, 'fg': GRAY},
], rx, CY + Inches(3.55), rw, Inches(1.2))

label(s, "Йменування варіантів:   squat.json  |  squat_front.json  |  squat_lowbar.json",
      CL, Inches(7.0), CW, Inches(0.42), sz=15, fg=DIM, mono=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Motion Signature
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Motion Signature — 99-вимірний вектор ознак",
        "Обчислюється на canonical rep (усереднений повтор), не на всій сесії")

lw = Inches(7.2)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

mk_table(s,
    ["Компонент", "Розмір", "Що кодує"],
    [("mean",                "9",   "Типова постура + проекція камери"),
     ("std",                 "9",   "Амплітуда руху кожного суглоба"),
     ("p10 / p90",           "18",  "Робастні екстремуми (глибина присіду)"),
     ("skewness",            "9",   "Час у нижній vs верхній точці"),
     ("mean_abs_velocity",   "9",   "Середня швидкість (deg/s)"),
     ("velocity_asymmetry",  "9",   "Частка кадрів де суглоб розгинається"),
     ("correlations (upper triangle)", "36", "Попарна кореляція суглобів"),
     ("TOTAL",               "99",  "float32")],
    CL, CY, lw, Inches(4.5),
    col_widths=[lw*0.40, lw*0.12, lw*0.48],
    hsz=14, rsz=13)

label(s, "Метрика схожості:", rx, CY, rw, Inches(0.42), sz=16, bold=True, fg=LBLUE)
code_block(s,
    "sim(u, r) =\n"
    "  (u · r) / (‖u‖ · ‖r‖)\n\n"
    "Косинусна відстань:\n"
    "  d = 1 - sim",
    rx, CY + Inches(0.48), rw, Inches(2.0), sz=14)

multiline(s, [
    {'text': 'Чому косинусна?', 'sz': 15, 'fg': WHITE, 'bold': True},
    {'text': 'Масштабно-інваріантна:', 'sz': 14, 'fg': GRAY},
    {'text': 'кут коліна ~120° не домінує над нахилом тулуба ~10°', 'sz': 14, 'fg': GRAY},
], rx, CY + Inches(2.65), rw, Inches(1.5))

multiline(s, [
    {'text': 'Навіщо canonical rep?', 'sz': 15, 'fg': WHITE, 'bold': True, 'sp_b': 8},
    {'text': 'Signature на всій сесії спотворена паузами між підходами. Canonical rep — один чистий цикл без паразитних фреймів.', 'sz': 14, 'fg': GRAY},
], CL, CY + Inches(4.65), CW, Inches(1.6))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Reference matching
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Автоматичний вибір найближчого референсу")

lw = Inches(6.3)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

code_block(s,
    "def find_best_match(user_ts, exercise):\n"
    "    # 1. Завантажити всі варіанти\n"
    "    candidates = load_family(exercise)\n"
    "    #  squat.json, squat_front.json, squat_lowbar.json\n\n"
    "    if len(candidates) == 1:\n"
    "        return candidates[0]\n\n"
    "    # 2. Обчислити signature користувача\n"
    "    #    автосегментація → canonical rep → 99-dim вектор\n"
    "    user_sig = MotionSignature.from_timeseries(user_ts)\n\n"
    "    # 3. Знайти найближчий (мінімальна косинусна відстань)\n"
    "    return min(candidates,\n"
    "               key=lambda r: user_sig.distance(r.signature))",
    CL, CY, lw, Inches(4.0), sz=13)

label(s, "Лог виводу:", CL, CY + Inches(4.15), lw, Inches(0.4), sz=15, bold=True, fg=LBLUE)
code_block(s,
    "[ReferenceStore] 3 variants for 'squat'\n"
    "  ✓ best:    'squat'        dist=0.031\n"
    "  ✗ skipped: 'squat_lowbar' dist=0.187\n"
    "  ✗ skipped: 'squat_front'  dist=0.312",
    CL, CY + Inches(4.6), lw, Inches(1.15), sz=13)

multiline(s, [
    {'text': 'Як signature кодує ракурс?', 'sz': 16, 'fg': WHITE, 'bold': True, 'sp_a': 6},
    {'text': 'Той самий суглоб проектується по-різному з різних точок зйомки:', 'sz': 15, 'fg': GRAY, 'sp_a': 4},
    {'text': '  Збоку:    коліно  170°→90°  (Δ≈80°)', 'sz': 14, 'fg': LBLUE, 'mono': True, 'sp_a': 2},
    {'text': '  Спереду:  коліно  170°→130° (Δ≈40°)', 'sz': 14, 'fg': LBLUE, 'mono': True, 'sp_a': 6},
    {'text': 'Mean, std, p10/p90 — всі вони різні.', 'sz': 15, 'fg': GRAY, 'sp_a': 4},
    {'text': 'Косинусна відстань автоматично знаходить найближчий ракурс.', 'sz': 15, 'fg': GRAY},
], rx, CY, rw, Inches(5.75))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Rep segmentation
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Сегментація повторів")

steps_seg = [
    ("Крок 1  —  PCA-сигнал",
     "Замість одного суглоба — 1-а головна компонента всіх 9 ознак.\n"
     "PCA об'єднує корельовані рухи і усереднює шум окремих датчиків.\n"
     "Орієнтація: позитивно корельований з найбільш амплітудною ознакою."),
    ("Крок 2  —  Multi-scale згладжування",
     "Три розміри вікна: ×1, ×2, ×4 від базового (Savitzky-Golay).\n"
     "Грубіше згладжування прибирає внутрішньоповторні піки.\n"
     "Запускаємо sweep на всіх трьох — беремо найкращий результат."),
    ("Крок 3  —  Adaptive prominence sweep",
     "Порогова умова: max(min_prominence, ratio × amplitude).\n"
     "Перевага позитивним пікам (вершини = стоячи = природна межа rep).\n"
     "Голова/хвіст включаються лише якщо ≥ max(min_frames, 0.6×median)."),
]

bh = Inches(1.8)
gap = Inches(0.15)
for i, (title, body) in enumerate(steps_seg):
    by = CY + i * (bh + gap)
    box(s, CL, by, CW, bh, fill=BOXBG, line=BLUE)
    box(s, CL, by, Inches(0.55), bh, fill=TH)
    label(s, str(i+1), CL, by, Inches(0.55), bh,
          sz=26, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)
    label(s, title, CL + Inches(0.65), by + Inches(0.1), CW - Inches(0.8), Inches(0.5),
          sz=17, bold=True, fg=LBLUE)
    multiline(s, [{'text': body, 'sz': 15, 'fg': GRAY}],
              CL + Inches(0.65), by + Inches(0.55), CW - Inches(0.8), Inches(1.1))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 15 — Segmentation math
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Математика сегментації")

lw = Inches(6.1)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

# Left: CV + scoring
label(s, "Коефіцієнт варіації (CV)", CL, CY, lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
multiline(s, [
    {'text': 'Справжні повтори займають приблизно однаковий час → рівні інтервали → мінімальний CV.', 'sz': 15, 'fg': GRAY},
], CL, CY + Inches(0.5), lw, Inches(0.8))
code_block(s, "CV = σ(intervals) / μ(intervals)", CL, CY + Inches(1.35), lw, Inches(0.65), sz=16)

label(s, "Функція оцінки кандидата", CL, CY + Inches(2.15), lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s, "score = CV + 0.002 × N_peaks", CL, CY + Inches(2.65), lw, Inches(0.65), sz=16)
multiline(s, [
    {'text': 'При рівному CV → перевага меншій кількості піків (Occam\'s Razor).', 'sz': 15, 'fg': GRAY},
], CL, CY + Inches(3.4), lw, Inches(0.65))

label(s, "PCA-сигнал", CL, CY + Inches(4.15), lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s,
    "_, _, Vt = np.linalg.svd(centered_data)\n"
    "pc1 = centered_data @ Vt[0]   # (T,) перша г.к.",
    CL, CY + Inches(4.65), lw, Inches(1.1), sz=13)

# Right: +1 fix
box(s, rx, CY, rw, Inches(5.75), fill=BOXBG, line=RED, lw=1.0)
label(s, "Захист від помилки +1 повтор", rx + Inches(0.2), CY + Inches(0.1),
      rw - Inches(0.4), Inches(0.5), sz=15, bold=True, fg=RED)
multiline(s, [
    {'text': 'Проблема:', 'sz': 14, 'fg': WHITE, 'bold': True, 'sp_a': 4},
    {'text': 'Якщо піки потрапляють у нижню точку руху, і голова, і хвіст сегменту проходять поріг min_frames → N+1 повтор замість N.', 'sz': 13, 'fg': GRAY, 'sp_a': 8},
    {'text': 'Рішення 1:', 'sz': 14, 'fg': WHITE, 'bold': True, 'sp_a': 4},
    {'text': 'Шукаємо позитивні піки першими. Вершини = стоячи = природна межа повтору. Голова і хвіст стають короткими.', 'sz': 13, 'fg': GRAY, 'sp_a': 8},
    {'text': 'Рішення 2:', 'sz': 14, 'fg': WHITE, 'bold': True, 'sp_a': 4},
    {'text': 'Поріг включення:', 'sz': 13, 'fg': GRAY, 'sp_a': 2},
], rx + Inches(0.2), CY + Inches(0.7), rw - Inches(0.4), Inches(4.0))
code_block(s,
    "threshold = max(min_frames,\n"
    "    int(0.6 × median_interior))",
    rx + Inches(0.1), CY + Inches(4.7), rw - Inches(0.2), Inches(0.9), sz=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 16 — DTW comparison
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "DTW порівняння з референсом")

lw = Inches(6.0)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

label(s, "Чому не Euclidean?", CL, CY, lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s,
    "Euclidean:  потребує однакову довжину і\n"
    "            кадрове вирівнювання — нереально.\n\n"
    "DTW:  знаходить оптимальне вирівнювання.\n"
    "  User  [A, B, B, C, D, E]  (повільніший підхід)\n"
    "  Ref   [a,    b, c, d, e]\n"
    "  Path: (1,1)(2,2)(3,2)(4,3)(5,4)(6,5)",
    CL, CY + Inches(0.5), lw, Inches(2.5), sz=13)

label(s, "FastDTW + адаптивний радіус", CL, CY + Inches(3.1), lw, Inches(0.45),
      sz=17, bold=True, fg=LBLUE)
code_block(s,
    "Точний DTW:  O(N²) — нереально для 300 кадрів\n"
    "FastDTW:     O(N)  — Sakoe-Chiba band\n\n"
    "radius = max(r, int(0.15 × max(T_q, T_r)))",
    CL, CY + Inches(3.6), lw, Inches(1.85), sz=13)

label(s, "Вагові коефіцієнти ознак", rx, CY, rw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
mk_table(s,
    ["Ознака", "Вага"],
    [("torso_lean",       "2.5"),
     ("left_knee_angle",  "2.0"),
     ("right_knee_angle", "2.0"),
     ("left_hip_angle",   "1.5"),
     ("right_hip_angle",  "1.5"),
     ("left_elbow_angle", "1.0"),
     ("right_elbow_angle","1.0"),
     ("left_shoulder_angle","1.0"),
     ("right_shoulder_angle","1.0")],
    rx, CY + Inches(0.5), rw, Inches(4.95),
    col_widths=[rw*0.72, rw*0.28],
    hsz=13, rsz=12)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 17 — Tolerance band
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Толерантна зона референсу",
        "Деякі відхилення — це природна варіативність, а не помилка")

lw = Inches(6.2)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

multiline(s, [
    {'text': 'Якщо референс містить N > 1 повторів:', 'sz': 17, 'fg': WHITE, 'bold': True, 'sp_a': 6},
], CL, CY, lw, Inches(0.55))
code_block(s,
    "# Усі повтори → медіанна довжина → стек\n"
    "stacked = np.stack([rep.resample(T).data\n"
    "                    for rep in reps])  # (N, T, F)\n\n"
    "canonical  = stacked.mean(axis=0)  # (T, F)\n"
    "tolerance  = stacked.std(axis=0)   # (T, F)",
    CL, CY + Inches(0.6), lw, Inches(2.1), sz=13)

label(s, "Застосування при порівнянні:", CL, CY + Inches(2.8), lw, Inches(0.45),
      sz=17, bold=True, fg=LBLUE)
code_block(s,
    "# DTW path: q_idx → r_idx\n"
    "tol = tolerance[r_idx]          # (F,)\n"
    "abs_dev = |raw_dev|\n"
    "effective = sign(raw) × max(0, abs_dev − tol)",
    CL, CY + Inches(3.3), lw, Inches(1.75), sz=13)

label(s, "Приклад (один кадр):", CL, CY + Inches(5.15), lw, Inches(0.4), sz=15, bold=True, fg=LBLUE)
code_block(s,
    "Tolerance: [ 3°,  5°,  2°, ... ]\n"
    "Raw dev:   [ 7°,  4°,  1°, ... ]\n"
    "Effective: [ 4°,  0°,  0°, ... ]  ← лише понад межу",
    CL, CY + Inches(5.6), lw, Inches(1.15), sz=13)

multiline(s, [
    {'text': 'Переваги:', 'sz': 16, 'fg': WHITE, 'bold': True, 'sp_a': 6},
    {'text': '  •  Не штрафуємо природну варіативність еталону', 'sz': 15, 'fg': GRAY, 'sp_a': 4},
    {'text': '  •  Відхилення всередині σ еталону = правильно', 'sz': 15, 'fg': GRAY, 'sp_a': 4},
    {'text': "  •  Більше ref-повторів → ширша зона → м'якший штраф", 'sz': 15, 'fg': GRAY, 'sp_a': 4},
    {'text': '', 'sz': 6},
    {'text': 'При одному ref-повторі:', 'sz': 15, 'fg': DIM, 'sp_a': 4},
    {'text': '  tolerance = 0 скрізь → кожне відхилення враховується', 'sz': 14, 'fg': DIM},
], rx, CY, rw, Inches(5.75))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 18 — Scoring math
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Оцінка техніки: математика")

lw = Inches(6.8)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

label(s, "Формула оцінки кадру:", CL, CY, lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s,
    "frame_score(t) = 100 × exp(−k × mean_|dev|(t))\n\n"
    "де:\n"
    "  mean_|dev|(t) = mean(|effective_dev(t, f)|\n"
    "                       for f in features)\n"
    "  k = sensitivity = 0.05  (конфіг.)",
    CL, CY + Inches(0.5), lw, Inches(2.35), sz=14)

label(s, "Оцінка повтору і сесії:", CL, CY + Inches(2.95), lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s,
    "rep_score     = mean(frame_score(t)  for t in rep)\n"
    "overall_score = mean(rep_score(i)    for i in reps)",
    CL, CY + Inches(3.45), lw, Inches(1.2), sz=14)

label(s, "Консистентність:", CL, CY + Inches(4.75), lw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
code_block(s,
    "consistency = std(rep_scores)   # нижче = стабільніше",
    CL, CY + Inches(5.25), lw, Inches(0.7), sz=14)

label(s, "Градуювання:", rx, CY, rw, Inches(0.45), sz=17, bold=True, fg=LBLUE)
mk_table(s,
    ["Бали", "Грейд", "Значення"],
    [("90–100", "A", "Відмінно"),
     ("80–89",  "B", "Добре"),
     ("70–79",  "C", "Непогано"),
     ("60–69",  "D", "Слабко"),
     ("< 60",   "F", "Серйозні проблеми")],
    rx, CY + Inches(0.5), rw, Inches(3.2),
    col_widths=[rw*0.30, rw*0.20, rw*0.50],
    hsz=14, rsz=13, center_all=False)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 19 — Error classification
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Класифікація помилок")

lw = Inches(6.0)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

# Left: levels
multiline(s, [
    {'text': 'Два рівні', 'sz': 17, 'fg': WHITE, 'bold': True, 'sp_a': 6},
], CL, CY, lw, Inches(0.5))

for i, (level, threshold, color, desc) in enumerate([
    ("CRITICAL",  "> 12°", RED,    "Ризик травми — негайне виправлення"),
    ("TECHNICAL", "> 8°",  YELLOW, "Ефективність руху — варто покращити"),
]):
    by = CY + Inches(0.6) + i * Inches(1.35)
    box(s, CL, by, lw, Inches(1.2), fill=BOXBG, line=color, lw=1.2)
    label(s, level, CL + Inches(0.2), by + Inches(0.1), Inches(2.2), Inches(0.55),
          sz=20, bold=True, fg=color)
    label(s, threshold, CL + Inches(2.5), by + Inches(0.1), Inches(1.2), Inches(0.55),
          sz=18, bold=True, fg=WHITE, mono=True)
    label(s, desc, CL + Inches(0.2), by + Inches(0.65), lw - Inches(0.4), Inches(0.5),
          sz=15, fg=GRAY)

label(s, "Критичні суглоби (вищий DTW-вага):", CL, CY + Inches(3.1), lw, Inches(0.45),
      sz=16, bold=True, fg=LBLUE)
code_block(s,
    "CRITICAL_FEATURES = {\n"
    "    'left_knee_angle', 'right_knee_angle', 'torso_lean'\n"
    "}",
    CL, CY + Inches(3.6), lw, Inches(1.15), sz=13)

label(s, "Направлення помилки:", CL, CY + Inches(4.9), lw, Inches(0.45),
      sz=16, bold=True, fg=LBLUE)
code_block(s,
    "+ відхилення → суглоб перерозігнутий (занадто відкритий)\n"
    "- відхилення → суглоб недорозігнутий (мала глибина)",
    CL, CY + Inches(5.4), lw, Inches(1.1), sz=13)

# Right: example
label(s, "Приклад виводу:", rx, CY, rw, Inches(0.45), sz=16, bold=True, fg=LBLUE)
code_block(s,
    "left_knee_angle:\n"
    "  avg_dev = +37°\n"
    "  → CRITICAL\n"
    "    (гіперекстензія)\n\n"
    "torso_lean:\n"
    "  avg_dev = +16°\n"
    "  → CRITICAL\n"
    "    (надмірний нахил)\n\n"
    "left_hip_angle:\n"
    "  avg_dev = +10°\n"
    "  → TECHNICAL",
    rx, CY + Inches(0.5), rw, Inches(5.25), sz=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 20 — Feedback system
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Система фідбеку")

lw = Inches(6.6)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

label(s, "Rule-based генерація:", CL, CY, lw, Inches(0.45), sz=16, bold=True, fg=LBLUE)
code_block(s,
    "CUES = {\n"
    "  ('left_knee_angle', POSITIVE): (\n"
    "    'Don\\'t lock out your left knee',\n"
    "    'Keep a slight bend at the top.'\n"
    "  ),\n"
    "  ('torso_lean', POSITIVE): (\n"
    "    'Keep your chest up',\n"
    "    'Drive through your heels, brace core.'\n"
    "  ),\n"
    "  ...\n"
    "}",
    CL, CY + Inches(0.5), lw, Inches(2.95), sz=13)

label(s, "Пріоритизація:", CL, CY + Inches(3.6), lw, Inches(0.45), sz=16, bold=True, fg=LBLUE)
code_block(s,
    "1. CRITICAL помилки (> 12°)\n"
    "2. TECHNICAL помилки (> 8°)\n"
    "3. В межах кожного рівня: сортування за |відхиленням|",
    CL, CY + Inches(4.1), lw, Inches(1.2), sz=14)

multiline(s, [
    {'text': 'Що включає звіт:', 'sz': 16, 'fg': WHITE, 'bold': True, 'sp_a': 6},
    {'text': '  •  Загальна оцінка + грейд (A–F)', 'sz': 14, 'fg': GRAY, 'sp_a': 4},
    {'text': '  •  Per-rep розбивка (найкращий / найгірший)', 'sz': 14, 'fg': GRAY, 'sp_a': 4},
    {'text': '  •  Топ-3 пріоритетних виправлення', 'sz': 14, 'fg': GRAY, 'sp_a': 4},
    {'text': '  •  Зведена порада для сесії', 'sz': 14, 'fg': GRAY, 'sp_a': 4},
    {'text': '  •  Консистентність (std повторів)', 'sz': 14, 'fg': GRAY},
], rx, CY, rw, Inches(5.75))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 21 — Video annotation
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Анотація відео")

lw = Inches(5.8)
rw = CW - lw - Inches(0.3)
rx = CL + lw + Inches(0.3)

label(s, "Кольорове кодування скелету:", CL, CY, lw, Inches(0.45), sz=16, bold=True, fg=LBLUE)

for i, (threshold, color, label_text, c) in enumerate([
    ("< 8°",    GREEN,  "Норма",              GREEN),
    ("8°–12°",  YELLOW, "Технічна помилка",   YELLOW),
    ("> 12°",   RED,    "Критична (ризик!)",  RED),
]):
    by = CY + Inches(0.55) + i * Inches(0.7)
    box(s, CL, by, Inches(0.55), Inches(0.55), fill=c)
    label(s, threshold, CL + Inches(0.65), by, Inches(1.0), Inches(0.55),
          sz=15, fg=WHITE, mono=True)
    label(s, label_text, CL + Inches(1.75), by, Inches(2.8), Inches(0.55),
          sz=15, fg=GRAY)

label(s, "Елементи кадру:", CL, CY + Inches(2.65), lw, Inches(0.45), sz=16, bold=True, fg=LBLUE)
code_block(s,
    "┌─────────────────────────────────────┐\n"
    "│  Squat              [████████░░] 82 │ ← шкала\n"
    "│                                     │\n"
    "│     * (плече)                       │\n"
    "│    /|\\                              │ ← скелет\n"
    "│   / | \\                             │\n"
    "│  /  |  \\                            │\n"
    "│ * (коліно — червоне)                │\n"
    "│                                     │\n"
    "│ ⚠ Bend your left knee more         │ ← банер\n"
    "└─────────────────────────────────────┘",
    CL, CY + Inches(3.15), lw, Inches(3.6), sz=12)

multiline(s, [
    {'text': 'Другий прохід:', 'sz': 16, 'fg': WHITE, 'bold': True, 'sp_a': 6},
    {'text': 'Аналіз зберігає кути суглобів (float32), а не піксельні координати.', 'sz': 14, 'fg': GRAY, 'sp_a': 6},
    {'text': 'Рендеринг повторно запускає PoseDetector у static_image_mode → отримує (x, y) у пікселях → малює.', 'sz': 14, 'fg': GRAY, 'sp_a': 8},
    {'text': 'Перевага: аналіз і рендеринг незалежні; аналіз можна запустити без відеовиводу.', 'sz': 14, 'fg': DIM},
], rx, CY, rw, Inches(5.75))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 22 — Limitations
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Обмеження та слабкі місця")

limitations = [
    ("Чутливість до кута зйомки",
     "Збоку: коліно 170°→90° (Δ≈80°). Спереду: 170°→130° (Δ≈40°).\n"
     "Рішення: Motion Signature знаходить найближчий ракурс автоматично.\n"
     "Потрібно: калібрування або 3D-реконструкція пози.",
     YELLOW),
    ("Окклюзія суглобів",
     "MediaPipe втрачає точки при поганому освітленні або ракурсі.\n"
     "Кадри з низькою visibility потрапляють у TimeSeries як шум.",
     YELLOW),
    ("Тільки 2D кути",
     "Проекція 3D руху на 2D площину → геометричні спотворення.\n"
     "Ідеально: стереокамера або IMU-датчики.",
     YELLOW),
    ("Rule-based фідбек",
     "Статичні правила. Не враховує контекст: стомленість vs структурний дефект.",
     YELLOW),
    ("Якість сегментації",
     "При нерівномірному темпі або нечіткому поверненні в стартову позицію\nможливі зайві або пропущені повтори.",
     YELLOW),
]

bh = Inches(1.12)
gap = Inches(0.08)
for i, (title, body, color) in enumerate(limitations):
    by = CY + i * (bh + gap)
    box(s, CL, by, CW, bh, fill=BOXBG, line=color, lw=0.75)
    label(s, f"  {i+1}.  {title}", CL, by + Inches(0.08), Inches(5.5), Inches(0.45),
          sz=15, bold=True, fg=color)
    multiline(s, [{'text': f"       {body}", 'sz': 13, 'fg': GRAY}],
              CL, by + Inches(0.52), CW, Inches(0.55))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 23 — Development plans
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Плани розвитку")

horizons = [
    ("3–6 місяців",
     ["Мобільний додаток (React Native + MediaPipe Mobile SDK)",
      "Push-сповіщення після тренування",
      "Хмарний бекенд — прогрес, синхронізація",
      "Веб-дашборд для тренерів"],
     BLUE),
    ("6–12 місяців",
     ["B2B інтеграція зі спортзалами (камери на стійках)",
      "CRM-інтеграція для тренерів",
      "Бібліотека 50+ вправ від сертифікованих тренерів",
      "LLM-фідбек — контекстний, адаптивний"],
     LBLUE),
    ("12+ місяців",
     ["3D pose estimation (стереокамера / мультикамера)",
      "AI-програма прогресії тренувань",
      "API для партнерів — інтеграція у фітнес-платформи",
      "Аналіз травмонебезпечних рухів (превентивний режим)"],
     GREEN),
]

bw = Inches(4.0)
gap = Inches(0.27)
for i, (period, items, color) in enumerate(horizons):
    bx = CL + i * (bw + gap)
    box(s, bx, CY, bw, Inches(5.7), fill=BOXBG, line=color, lw=1.0)
    box(s, bx, CY, bw, Inches(0.62), fill=TH)
    label(s, period, bx + Inches(0.2), CY + Inches(0.08), bw - Inches(0.4), Inches(0.5),
          sz=17, bold=True, fg=color, align=PP_ALIGN.CENTER)
    lines = [{'text': f"  •  {item}", 'sz': 14, 'fg': GRAY, 'sp_a': 8} for item in items]
    multiline(s, lines, bx + Inches(0.15), CY + Inches(0.75), bw - Inches(0.3), Inches(4.75))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 24 — Comparison of approaches
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Порівняння підходів до аналізу руху")

mk_table(s,
    ["Підхід", "Переваги", "Недоліки"],
    [("Кути + DTW  (наш підхід)",
      "Інваріантний до зросту/камери; інтерпретований",
      "2D; потребує схожий ракурс"),
     ("Сирі координати (x,y)",
      "Просто реалізувати",
      "Залежить від дистанції та зросту"),
     ("Pixel-level відео (CNN)",
      "Не потребує детекції пози",
      "Чорний ящик; важко пояснити"),
     ("IMU-сенсори",
      "Точно; справжнє 3D",
      "Спеціальне обладнання; незручно"),
     ("GPT/Vision LLM",
      "Гнучкий, людяний фідбек",
      "Повільно; дорого; нестабільно")],
    CL, CY, CW, Inches(4.5),
    col_widths=[CW*0.28, CW*0.38, CW*0.34],
    hsz=15, rsz=14)

box(s, CL, CY + Inches(4.65), CW, Inches(1.6), fill=TH, line=BLUE)
multiline(s, [
    {'text': 'Ключова перевага нашого підходу:', 'sz': 16, 'fg': WHITE, 'bold': True, 'sp_a': 6},
    {'text': 'Результат пояснюваний — тренер бачить відхилення по кожному суглобу і розуміє, чому система дала таку оцінку.', 'sz': 16, 'fg': GRAY},
], CL + Inches(0.3), CY + Inches(4.75), CW - Inches(0.6), Inches(1.4))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 25 — Summary
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Підсумок")

cols_sum = [
    ("✅  Що зроблено",
     ["Повний прототип: відео → фідбек",
      "Незалежний від кількості повторів",
      "Автоматичний вибір найближчого референсу",
      "Готовий до інтеграції в мобільний додаток",
      "CLI: record / analyze / list / demo"],
     GREEN),
    ("🎯  Унікальне позиціювання",
     ["Тренер записує власний еталон вправи",
      "Система адаптується під його методику",
      "Масштабування без залучення тренера",
      "Пояснюваний AI (не чорний ящик)",
      "JSON-сховище — Git-friendly, портативне"],
     BLUE),
    ("📈  Наступний крок",
     ["Мобільний MVP",
      "3 пілотних спортзали для валідації",
      "Збір реального датасету помилок",
      "LLM-фідбек замість rule-based",
      "3D pose estimation"],
     LBLUE),
]

bw = Inches(4.0)
gap = Inches(0.27)
for i, (title, items, color) in enumerate(cols_sum):
    bx = CL + i * (bw + gap)
    box(s, bx, CY, bw, Inches(5.7), fill=BOXBG, line=color, lw=1.0)
    box(s, bx, CY, bw, Inches(0.65), fill=TH)
    label(s, title, bx + Inches(0.18), CY + Inches(0.1), bw - Inches(0.36), Inches(0.52),
          sz=16, bold=True, fg=color)
    lines = [{'text': f"  •  {item}", 'sz': 14, 'fg': GRAY, 'sp_a': 8} for item in items]
    multiline(s, lines, bx + Inches(0.15), CY + Inches(0.78), bw - Inches(0.3), Inches(4.65))

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 26 — Thank you
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
box(s, 0, H - Inches(0.55), W, Inches(0.55), fill=BAR)
box(s, 0, H - Inches(0.55), W, Inches(0.07), fill=BLUE)

label(s, "Дякуємо за увагу!",
      Inches(0.8), Inches(1.0), Inches(11.8), Inches(1.5),
      sz=50, bold=True, fg=WHITE, align=PP_ALIGN.CENTER)

label(s, "AI Fitness Coach",
      Inches(0.8), Inches(2.5), Inches(11.8), Inches(0.7),
      sz=28, fg=LBLUE, align=PP_ALIGN.CENTER)

box(s, Inches(3.8), Inches(3.4), Inches(5.8), Inches(2.5), fill=BOXBG, line=BLUE)
multiline(s, [
    {'text': '  Email:    denis.stepchin123@gmail.com', 'sz': 18, 'fg': GRAY, 'mono': True, 'sp_a': 8},
    {'text': '  GitHub:   github.com / ai-fitness-coach', 'sz': 18, 'fg': GRAY, 'mono': True, 'sp_a': 8},
    {'text': '  Mobile:   [в розробці]', 'sz': 18, 'fg': DIM, 'mono': True},
], Inches(3.9), Inches(3.5), Inches(5.6), Inches(2.2))

label(s, '"Ваш персональний тренер, доступний 24/7 — у кишені"',
      Inches(1.5), Inches(6.2), Inches(10.4), Inches(0.7),
      sz=18, italic=True, fg=DIM, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 27 — Appendix A: Project structure
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Додаток A  —  Структура проекту")

struct = (
    "ai_trainer/\n"
    "├── main.py                      CLI: record / analyze / list / demo\n"
    "├── src/\n"
    "│   ├── pose/\n"
    "│   │   ├── detector.py          MediaPipe Tasks API wrapper\n"
    "│   │   └── keypoints.py         PoseFrame, PoseLandmark\n"
    "│   ├── features/\n"
    "│   │   ├── extractor.py         9 кутових ознак  ← FEATURE_NAMES контракт\n"
    "│   │   └── normalizer.py        Hip-center + torso-scale\n"
    "│   ├── timeseries/\n"
    "│   │   ├── builder.py           TimeSeries(T×F), smoothing, resampling\n"
    "│   │   └── segmentation.py      RepSegmenter: PCA + multi-scale peaks\n"
    "│   ├── comparison/\n"
    "│   │   ├── dtw.py               DTWComparator, FastDTW, weights\n"
    "│   │   └── rep_comparator.py    RepComparator, tolerance band\n"
    "│   ├── evaluation/\n"
    "│   │   ├── scorer.py            ScoringEngine: exponential decay\n"
    "│   │   └── classifier.py        ErrorClassifier: Critical / Technical\n"
    "│   ├── feedback/\n"
    "│   │   └── generator.py         FeedbackGenerator: rule-based cues\n"
    "│   ├── reference/\n"
    "│   │   ├── store.py             ReferenceStore: CRUD + find_best_match\n"
    "│   │   └── signature.py         MotionSignature: 99-dim embedding\n"
    "│   ├── visualization/\n"
    "│   │   ├── overlay.py           FrameAnnotator\n"
    "│   │   └── renderer.py          VideoRenderer (2nd pass)\n"
    "│   └── pipeline.py              AnalysisPipeline — orchestrator\n"
    "└── references/                  JSON reference files (gitignored)"
)
code_block(s, struct, CL, CY, CW, Inches(5.88), sz=11.5)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 28 — Appendix B: System invariants
# ═══════════════════════════════════════════════════════════════════════════════
s = new_slide()
top_bar(s, "Додаток B  —  Ключові інваріанти системи")

invariants = [
    ("FEATURE_NAMES — глобальний контракт",
     "Порядок у extractor.py визначає порядок стовпців у кожному масиві (T×F). Зміна = зламані всі збережені референси та весь downstream.",
     RED),
    ("Сегментація перед DTW",
     "RepComparator завжди сегментує обидві послідовності і порівнює 1-проти-1. Ніколи не викликати DTWComparator безпосередньо на повній сесії.",
     YELLOW),
    ("combined_dtw — синтетичний об'єкт",
     "RepComparisonResult.combined_dtw — конкатенація per-rep результатів з ідентичним alignment path. Не використовувати distance або path для per-rep аналізу.",
     YELLOW),
    ("MediaPipe Tasks API — тільки",
     "mediapipe.tasks.python.vision.PoseLandmarker. Ніколи mp.solutions.* — відсутній на Windows у MediaPipe ≥ 0.10.",
     BLUE),
    ("Візуалізація — другий прохід",
     "Аналіз зберігає кути, а не піксельні координати. Рендеринг повторно запускає детектор у static_image_mode.",
     BLUE),
]

bh = Inches(1.1)
gap = Inches(0.09)
for i, (title, body, color) in enumerate(invariants):
    by = CY + i * (bh + gap)
    box(s, CL, by, CW, bh, fill=BOXBG, line=color, lw=0.75)
    label(s, f"  {i+1}.  {title}", CL, by + Inches(0.07), Inches(6.0), Inches(0.44),
          sz=15, bold=True, fg=color)
    multiline(s, [{'text': f"       {body}", 'sz': 13, 'fg': GRAY}],
              CL, by + Inches(0.52), CW, Inches(0.52))

# ── Save ──────────────────────────────────────────────────────────────────────
out = "AI_Fitness_Coach.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
